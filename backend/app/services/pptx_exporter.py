from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


BRAND_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
BRAND_LIGHT = RGBColor(0xF8, 0xFA, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(tf, text: str, size: int, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or BRAND_BLUE


def generate_pptx(dashboard_data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s, BRAND_BLUE)

    tx = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.2))
    _add_text(tx.text_frame, dashboard_data.get('title', 'Dashboard'), 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tx2 = s.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(10), Inches(1.8))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    _add_text(tf2, dashboard_data.get('executive_summary', ''), 18, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

    tx3 = s.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.6))
    _add_text(tx3.text_frame, "Powered by Vizify", 12, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 2: KPIs ───────────────────────────────────────────────────────
    kpis = dashboard_data.get('kpis', [])
    if kpis:
        s = prs.slides.add_slide(blank)
        _set_bg(s, BRAND_LIGHT)

        hdr = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        _add_text(hdr.text_frame, "Key Performance Indicators", 28, bold=True, color=BRAND_BLUE)

        cols = min(len(kpis), 4)
        w = Inches(11.5 / cols)
        for i, kpi in enumerate(kpis[:8]):
            col = i % 4
            row = i // 4
            x = Inches(0.5) + col * w
            y = Inches(1.3) + row * Inches(2.8)
            box = s.shapes.add_textbox(x, y, w - Inches(0.2), Inches(2.5))
            tf = box.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = kpi.get('label', '')
            r1.font.size = Pt(13)
            r1.font.color.rgb = GRAY

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = kpi.get('value', '')
            r2.font.size = Pt(30)
            r2.font.bold = True
            r2.font.color.rgb = BRAND_BLUE

            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run()
            trend = kpi.get('trend', 'flat')
            arrow = '↑' if trend == 'up' else ('↓' if trend == 'down' else '→')
            r3.text = f"{arrow} {kpi.get('delta', '')}"
            r3.font.size = Pt(14)
            r3.font.bold = True
            r3.font.color.rgb = GREEN if trend == 'up' else (RED if trend == 'down' else GRAY)

    # ── Slide 3: Insights ───────────────────────────────────────────────────
    insights = dashboard_data.get('insights', [])
    if insights:
        s = prs.slides.add_slide(blank)
        _set_bg(s, WHITE)

        hdr = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        _add_text(hdr.text_frame, "Key Insights", 28, bold=True, color=BRAND_BLUE)

        tx = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, ins in enumerate(insights[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"• {ins}"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    # ── Slide 4: Recommendations ─────────────────────────────────────────────
    recs = dashboard_data.get('recommendations', [])
    if recs:
        s = prs.slides.add_slide(blank)
        _set_bg(s, BRAND_BLUE)

        hdr = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        _add_text(hdr.text_frame, "Recommendations", 28, bold=True, color=WHITE)

        tx = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, rec in enumerate(recs[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"{i+1}.  {rec}"
            run.font.size = Pt(16)
            run.font.color.rgb = WHITE

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
